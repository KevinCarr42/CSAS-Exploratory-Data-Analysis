"""Helper functions for text extraction from meeting documents."""
import os
from pathlib import Path
from docx import Document
from pptx import Presentation

def extract_docx_content(file_path):
    """Extract text from .docx file with structure information."""
    doc = Document(file_path)
    content = []

    for para in doc.paragraphs:
        if para.text.strip():
            element_type = 'heading' if para.style.name.startswith('Heading') else 'paragraph'
            content.append({
                'text': para.text,
                'element_type': element_type,
                'style': para.style.name
            })

    if doc.tables:
        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                row_data = [cell.text for cell in row.cells]
                content.append({
                    'text': ' | '.join(row_data),
                    'element_type': 'table_row',
                    'style': f'table_{table_idx}_row_{row_idx}'
                })

    return content

def extract_pptx_content(file_path):
    """Extract text from .pptx file with slide information."""
    prs = Presentation(file_path)
    content = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                content.append({
                    'text': shape.text,
                    'element_type': 'slide_text',
                    'style': f'slide_{slide_idx}_shape_{shape_idx}'
                })

            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    row_data = [cell.text for cell in row.cells]
                    content.append({
                        'text': ' | '.join(row_data),
                        'element_type': 'table_row',
                        'style': f'slide_{slide_idx}_table_row_{row_idx}'
                    })

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                content.append({
                    'text': notes_text,
                    'element_type': 'slide_notes',
                    'style': f'slide_{slide_idx}_notes'
                })

    return content
